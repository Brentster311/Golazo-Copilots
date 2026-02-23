"""Generate an executive-review PowerPoint deck for SFI Reporter."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Palette ──────────────────────────────────────────────────────────────
BLUE    = RGBColor(0x00, 0x78, 0xD4)   # Microsoft blue
DARK    = RGBColor(0x1B, 0x1B, 0x1B)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0x60, 0x60, 0x60)
LTGRAY  = RGBColor(0xD0, 0xD0, 0xD0)
ACCENT  = RGBColor(0x00, 0xB2, 0x94)   # Teal accent
ORANGE  = RGBColor(0xFF, 0x8C, 0x00)
BG_DARK = RGBColor(0x24, 0x24, 0x24)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height


# ── Helpers ──────────────────────────────────────────────────────────────
def _solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or WHITE
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # smaller corner radius
    shape.adjustments[0] = 0.05
    return shape


def _text_box(slide, left, top, width, height, text, font_size=14, bold=False,
              color=DARK, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _bullet_frame(slide, left, top, width, height, items, font_size=14,
                  color=DARK, bold_first=False, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = spacing
        if bold_first and i == 0:
            p.font.bold = True
    return txBox


def _section_header(slide, text, top=Inches(0.4)):
    _text_box(slide, Inches(0.7), top, Inches(11), Inches(0.6),
              text, font_size=28, bold=True, color=BLUE)


def _footer(slide, text="SFI Reporter  |  ACCIA Team  |  Confidential"):
    _text_box(slide, Inches(0.7), Inches(6.9), Inches(6), Inches(0.4),
              text, font_size=9, color=LTGRAY)


def _card(slide, left, top, width, height, title, bullets, title_color=BLUE):
    """Draw a rounded card with a title and bullet list."""
    _add_shape(slide, left, top, width, height, fill_color=RGBColor(0xF9, 0xF9, 0xF9), line_color=LTGRAY)
    _text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.35),
              title, font_size=15, bold=True, color=title_color)
    _bullet_frame(slide, left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
                  bullets, font_size=12, color=DARK)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_solid_bg(slide, BG_DARK)

_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
          "SFI Reporter", font_size=48, bold=True, color=WHITE)
_text_box(slide, Inches(1), Inches(2.9), Inches(11), Inches(0.7),
          "Desktop Application for SFI / QEI Action-Item Management",
          font_size=22, color=LTGRAY)
_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
          "ACCIA Team  ·  Executive Review  ·  February 2026",
          font_size=16, color=GRAY)

# Accent bar
slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.8), Inches(3), Pt(4)).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = BLUE
slide.shapes[-1].line.fill.background()


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & Value Proposition
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, WHITE)
_section_header(slide, "Problem & Value Proposition")
_footer(slide)

_card(slide, Inches(0.7), Inches(1.2), Inches(5.8), Inches(2.5),
      "The Problem", [
          "SFI / QEI action items are tracked in S360 portal",
          "No offline or desktop-native access for ICs or managers",
          "ETA updates require manual navigation per item",
          "No automated risk analysis or bulk operations",
      ], title_color=ORANGE)

_card(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.5),
      "Our Solution", [
          "Native desktop app — no browser required",
          "One-click refresh with local caching (fast loads)",
          "Bulk & individual ETA editing from any view",
          "AI-powered analysis via Azure OpenAI (🤖 Analyze)",
      ], title_color=ACCENT)

_card(slide, Inches(0.7), Inches(4.0), Inches(11.9), Inches(2.5),
      "Key Outcomes", [
          "~70% reduction in time-to-action for ETA reviews vs. S360 web portal",
          "Managers get hierarchical org-tree view with drill-down to individual owners",
          "Standalone .exe distribution — zero Python install required for end users",
          "Reusable Python SDK (accia-s360) published to Azure Artifacts for org-wide adoption",
      ], title_color=BLUE)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Key Features (visual grid)
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, WHITE)
_section_header(slide, "Key Features")
_footer(slide)

features = [
    ("🔄  Refresh & Cache", [
        "Auto-detects alias via Azure CLI",
        "Color-coded cache-age indicator",
        "1-hour TTL, background refresh",
    ]),
    ("📋  ETA Management", [
        "Manual review (step-through items)",
        "Bulk auto-fix for invalid ETAs",
        "Per-view multi-select ETA update",
    ]),
    ("🔍  Smart Filtering", [
        "Clause-based filter builder (persisted)",
        "Service, Owner, Program, Due Date",
        "Column toggle with empty-indicators",
    ]),
    ("🤖  LLM Analysis", [
        "Right-click → Analyze with AI",
        "Mission · Steps · Resources · Risk",
        "Saved locally for repeat access",
    ]),
    ("👥  Org Hierarchy", [
        "Manager vs. IC view auto-detect",
        "Multi-level owner grouping tree",
        "Drill-down by service or owner",
    ]),
    ("⚙️  Resilience", [
        "KPI failure notification (orange ⚠️)",
        "One-click retry for failed KPIs",
        "SLA & ETA status columns in views",
    ]),
]

cols, rows = 3, 2
cw, ch = Inches(3.8), Inches(2.3)
x0, y0 = Inches(0.7), Inches(1.2)
gap = Inches(0.25)

for idx, (title, bullets) in enumerate(features):
    c, r = idx % cols, idx // cols
    left = x0 + c * (cw + gap)
    top  = y0 + r * (ch + gap)
    _card(slide, left, top, cw, ch, title, bullets)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture Overview
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, WHITE)
_section_header(slide, "Architecture Overview")
_footer(slide)

# Layers from top to bottom
layers = [
    ("SFI Reporter Desktop App (Tkinter)", "tk_app.py  ·  query_builder.py  ·  eta_logic.py  ·  llm_client.py",
     RGBColor(0xE8, 0xF0, 0xFE), BLUE),
    ("Data & Cache Layer", "data.py  ·  cache.py  ·  logging_config.py",
     RGBColor(0xE6, 0xF4, 0xEA), ACCENT),
    ("accia-s360 SDK", "client.py  ·  auth.py  ·  models.py  ·  endpoints/ (action_items · extended · discovery · graph)",
     RGBColor(0xFD, 0xF0, 0xE2), ORANGE),
    ("External APIs", "S360 REST API (v1 / v2)  ·  Microsoft Graph API  ·  Azure OpenAI",
     RGBColor(0xF5, 0xF5, 0xF5), GRAY),
]

box_width  = Inches(11)
box_height = Inches(1.15)
x0 = Inches(1.15)
y0 = Inches(1.3)
gap = Inches(0.2)

for i, (title, subtitle, bg, accent) in enumerate(layers):
    top = y0 + i * (box_height + gap)
    shape = _add_shape(slide, x0, top, box_width, box_height, fill_color=bg, line_color=accent)

    # accent left bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, top, Pt(6), box_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    _text_box(slide, x0 + Inches(0.3), top + Inches(0.15), Inches(10), Inches(0.4),
              title, font_size=16, bold=True, color=accent)
    _text_box(slide, x0 + Inches(0.3), top + Inches(0.6), Inches(10), Inches(0.4),
              subtitle, font_size=11, color=GRAY)

# Arrows between layers
for i in range(len(layers) - 1):
    top_of_arrow = y0 + (i + 1) * (box_height + gap) - gap
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.5), top_of_arrow, Inches(0.35), gap)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LTGRAY
    arrow.line.fill.background()


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — accia-s360 SDK Deep Dive
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, WHITE)
_section_header(slide, "accia-s360 — Reusable S360 Python SDK")
_footer(slide)

_card(slide, Inches(0.7), Inches(1.2), Inches(5.8), Inches(2.2),
      "What It Is", [
          "Standalone pip-installable Python package",
          "Published to Azure Artifacts (ACCIA feed)",
          "43 S360 + Graph API endpoints wrapped",
          "Typed models, structured exceptions, local cache",
      ])

_card(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.2),
      "Auth Chain", [
          "Step 1: AzureCliCredential (az login)",
          "Step 2: InteractiveBrowserCredential fallback",
          "Token scopes: S360 API + MS Graph",
          "Zero-config for users with Azure CLI installed",
      ])

_card(slide, Inches(0.7), Inches(3.7), Inches(3.85), Inches(2.7),
      "API Coverage (43 endpoints)", [
          "Action Items — ETA, owners, grid, summary",
          "KPIs — costs, metadata, target types",
          "Programs — objectives, waves (v2)",
          "Graph — user info, org tree, direct reports",
          "+ Reliability, Feature Flags, ADO, Lifecycle",
      ])

_card(slide, Inches(4.8), Inches(3.7), Inches(3.85), Inches(2.7),
      "Design Principles", [
          "azure-identity for credentials (no secrets)",
          "Retry with exponential backoff",
          "SHA-256 keyed JSON file cache",
          "Full type hints (Python 3.10+)",
          "Pytest suite: auth chain, payload, build, pkg",
      ])

_card(slide, Inches(8.9), Inches(3.7), Inches(3.7), Inches(2.7),
      "Adoption", [
          "SFI Reporter is primary consumer",
          "Available for any Python-based S360 tooling",
          "pip install accia-s360 --index-url ...",
          "from accia_s360 import S360Client",
          "3 lines to query your action items",
      ])


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Capability Map
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, WHITE)
_section_header(slide, "Capability Map (15 Registered Capabilities)")
_footer(slide)

cap_groups = [
    ("s360_client Core SDK", BLUE, [
        "s360-auth — Azure CLI credential mgmt",
        "s360-client-facade — orchestrates auth + cache + endpoints",
        "s360-models — typed dataclass models",
        "s360-exceptions — S360Error hierarchy",
        "s360-cache — JSON file cache with SHA-256 + TTL",
        "s360-action-items — ETA history, save ETAs",
        "s360-discovery — probe available endpoints",
        "s360-extended-endpoints — 40+ v1/v2 endpoints",
    ]),
    ("accia-s360 Package", ORANGE, [
        "accia-s360-auth — two-step auth (CLI → browser)",
        "accia-s360-client — redistributable S360Client",
        "accia-s360-endpoints — action items, extended, graph",
    ]),
    ("SFI Reporter App", ACCENT, [
        "reporter-tk-app — full desktop UI (Tkinter)",
        "reporter-data — parallel KPI grid loading",
        "reporter-cache — 1-hr TTL JSON cache",
        "reporter-eta-logic — propose/validate ETAs",
        "reporter-llm — Azure OpenAI analysis",
        "reporter-query-builder — filter clause UI",
        "reporter-logging — rotating file + console",
        "reporter-build — PyInstaller → .exe",
    ]),
]

x0-= Inches(0)
col_w = Inches(3.85)
gap = Inches(0.2)

for ci, (group_name, color, items) in enumerate(cap_groups):
    left = Inches(0.7) + ci * (col_w + gap)
    # Group header
    hdr = _add_shape(slide, left, Inches(1.2), col_w, Inches(0.45), fill_color=color)
    _text_box(slide, left + Inches(0.15), Inches(1.22), col_w - Inches(0.3), Inches(0.4),
              group_name, font_size=13, bold=True, color=WHITE)
    # Items
    _bullet_frame(slide, left + Inches(0.15), Inches(1.75), col_w - Inches(0.3), Inches(5),
                  items, font_size=10, color=DARK, spacing=Pt(4))


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Testing & Quality
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, WHITE)
_section_header(slide, "Testing & Quality")
_footer(slide)

_card(slide, Inches(0.7), Inches(1.2), Inches(5.8), Inches(5.0),
      "Test Suites", [
          "s360_client tests — auth, cache, client, discovery, endpoints",
          "accia-s360 tests — auth chain, ETA payload, build, package",
          "SFI Reporter tests — 15+ test files covering:",
          "  · cache, data, ETA logic, ETA UI",
          "  · query builder, detail modal colors",
          "  · per-KPI tests (SFI-023 thru SFI-035)",
          "  · Tk app integration tests",
          "",
          "All tests run via:  pytest tests/ -v",
          "Automated gate: tests must pass before build",
      ])

_card(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.3),
      "Distribution & Build", [
          "PyInstaller → single SFIReporter.exe",
          "Distributed as SFIReporter.zip (exe + README)",
          "Zero Python install needed for end users",
          "accia-s360 published to Azure Artifacts feed",
      ])

_card(slide, Inches(6.8), Inches(3.8), Inches(5.8), Inches(2.4),
      "Engineering Practices", [
          "Golazo Copilot workflow (9-role SDLC)",
          "Capability registry (capabilities.yaml)",
          "Impact analysis for every code change",
          "Type hints throughout (Python 3.10+)",
          "Rotating diagnostic logs (%TEMP%\\sfireporter\\)",
      ])


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Roadmap & Next Steps
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, WHITE)
_section_header(slide, "Roadmap & Next Steps")
_footer(slide)

_card(slide, Inches(0.7), Inches(1.2), Inches(3.85), Inches(5.0),
      "✅  Delivered", [
          "Desktop app with manager / IC views",
          "43-endpoint S360 SDK (accia-s360)",
          "Bulk & individual ETA management",
          "LLM-powered action-item analysis",
          "Smart filtering & column customization",
          "One-click retry for failed KPIs",
          "PyInstaller .exe distribution",
      ], title_color=ACCENT)

_card(slide, Inches(4.8), Inches(1.2), Inches(3.85), Inches(5.0),
      "🔜  Near-Term", [
          "Expand LLM analysis to batch mode",
          "Remediation playbook integration",
          "Notification/alert subscriptions",
          "Deeper org-hierarchy analytics",
          "Dashboard export (PDF/Excel)",
      ], title_color=BLUE)

_card(slide, Inches(8.9), Inches(1.2), Inches(3.7), Inches(5.0),
      "🔮  Future Vision", [
          "Web-based companion (ASP.NET / React)",
          "Teams integration for SFI alerts",
          "Cross-team benchmarking dashboards",
          "Auto-remediation via policy engine",
          "Self-service SDK onboarding portal",
      ], title_color=ORANGE)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Summary / Ask
# ══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
_solid_bg(slide, BG_DARK)

_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
          "Summary", font_size=36, bold=True, color=WHITE)

summary_items = [
    "SFI Reporter delivers a native desktop experience for SFI/QEI action-item management",
    "accia-s360 SDK provides a reusable, pip-installable Python client for 43 S360 + Graph APIs",
    "AI-powered analysis, smart filtering, and bulk ETA operations reduce manual effort significantly",
    "Standalone .exe distribution enables org-wide adoption with zero setup",
    "Comprehensive test coverage and Golazo Copilot SDLC ensure production-grade quality",
]

_bullet_frame(slide, Inches(1), Inches(2.5), Inches(11), Inches(3.5),
              summary_items, font_size=18, color=WHITE, spacing=Pt(14))

_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.6),
          "Questions?", font_size=28, bold=True, color=BLUE)

_text_box(slide, Inches(1), Inches(6.5), Inches(6), Inches(0.4),
          "ACCIA Team  ·  accia@microsoft.com", font_size=13, color=GRAY)


# ── Save ─────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "SFIReporter_Exec_Review.pptx")
prs.save(out_path)
print(f"✅  Saved: {out_path}")
